from pptx import Presentation
from pptx.exc import PackageNotFoundError
from typing import Dict, Any, List, Union
import os
import logging

# Setup logger for this module
logger = logging.getLogger(__name__)

def parse_pptx(file_path: Union[str, os.PathLike]) -> Dict[str, Any]:
    """
    Parses a .pptx file and extracts text and other relevant information.

    Args:
        file_path: Path to the .pptx file.

    Returns:
        A dictionary containing extracted data (e.g., text, speaker notes).
    """
    try:
        presentation = Presentation(file_path)
        extracted_data = {
            "text_from_slides": [],
            "speaker_notes": [],
            "metadata": {
                "author": presentation.core_properties.author,
                "last_modified_by": presentation.core_properties.last_modified_by,
                "revision": presentation.core_properties.revision,
                "title": presentation.core_properties.title,
                "subject": presentation.core_properties.subject,
                "keywords": presentation.core_properties.keywords,
                "category": presentation.core_properties.category,
                "comments": presentation.core_properties.comments,
                "num_slides": len(presentation.slides),
                "file_type": "pptx"
            }
        }

        total_text_extracted = 0
        slides_with_text = 0
        total_slides = len(presentation.slides)

        for i, slide in enumerate(presentation.slides):
            slide_text_parts = []
            slide_text_length = 0
            
            # Extract text from shapes
            for shape in slide.shapes:
                try:
                    # Handle text in shapes
                    if hasattr(shape, "text") and shape.text:
                        text = shape.text.strip()
                        if text:
                            slide_text_parts.append(f"Shape Text: {text}")
                            slide_text_length += len(text)
                            logger.debug(f"Slide {i+1}: Extracted {len(text)} chars from shape")
                    
                    # Handle text in tables
                    if shape.has_table:
                        table_text = []
                        for row_idx, row in enumerate(shape.table.rows):
                            row_text = []
                            for cell_idx, cell in enumerate(row.cells):
                                if cell.text.strip():
                                    cell_text = cell.text.strip()
                                    row_text.append(cell_text)
                                    slide_text_length += len(cell_text)
                            if row_text:
                                table_text.append(f"Row {row_idx+1}: {' | '.join(row_text)}")
                        if table_text:
                            slide_text_parts.append("Table Content:\n" + "\n".join(table_text))
                            logger.debug(f"Slide {i+1}: Extracted table with {len(table_text)} rows")
                    
                    # Handle text in placeholders
                    if shape.is_placeholder:
                        if shape.text.strip():
                            text = shape.text.strip()
                            slide_text_parts.append(f"Placeholder Text: {text}")
                            slide_text_length += len(text)
                            logger.debug(f"Slide {i+1}: Extracted {len(text)} chars from placeholder")
                    
                    # Handle text frames (common in PowerPoint)
                    if hasattr(shape, "text_frame") and shape.text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            if paragraph.text.strip():
                                text = paragraph.text.strip()
                                slide_text_parts.append(text)
                                slide_text_length += len(text)
                                logger.debug(f"Slide {i+1}: Extracted {len(text)} chars from text frame")
                
                except Exception as shape_error:
                    logger.warning(f"⚠️ Error processing shape in slide {i+1}: {str(shape_error)}")
                    continue
            
            # Combine all text from the slide
            slide_text = "\n".join(slide_text_parts)
            
            # Log extraction for this slide
            if slide_text:
                slides_with_text += 1
                total_text_extracted += slide_text_length
                logger.info(f"Slide {i+1}: Extracted {slide_text_length} characters")
                extracted_data["text_from_slides"].append({
                    "slide_number": i + 1,
                    "text": slide_text,
                    "text_length": slide_text_length
                })
            else:
                logger.warning(f"⚠️ No text extracted from slide {i+1}")

            # Extract speaker notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    extracted_data["speaker_notes"].append({
                        "slide_number": i + 1,
                        "notes": notes_text,
                        "notes_length": len(notes_text)
                    })
                    total_text_extracted += len(notes_text)
                    logger.info(f"Slide {i+1}: Extracted {len(notes_text)} characters from notes")
        
        # Log overall extraction quality
        if total_text_extracted == 0:
            logger.warning("⚠️ No text could be extracted from the PowerPoint file")
        else:
            logger.info(f"✅ Successfully extracted {total_text_extracted} characters from PowerPoint")
            logger.info(f"📊 Processed {total_slides} slides:")
            logger.info(f"   - {slides_with_text} slides contained text")
            logger.info(f"   - {len(extracted_data['speaker_notes'])} slides had speaker notes")
            if slides_with_text < total_slides:
                logger.warning(f"⚠️ Only {slides_with_text} out of {total_slides} slides contained extractable text")

        return extracted_data

    except PackageNotFoundError:
        logger.error(f"Error: '{file_path}' is not a valid PPTX file or is corrupted.", exc_info=True)
        return {"error": f"'{file_path}' is not a valid PPTX file or is corrupted."}
    except Exception as e:
        logger.error(f"Error parsing PPTX file {file_path}: {e}", exc_info=True)
        return {"error": f"An unexpected error occurred while parsing {file_path}: {str(e)}"}

if __name__ == '__main__':
    # Configure basic logging for direct script execution
    if not logging.getLogger().handlers:
        logging.basicConfig(level=logging.DEBUG, format='%(asctime)s - %(levelname)s - %(message)s')

    example_file = "example.pptx"
    try:
        # Attempt to create a dummy pptx for testing if it doesn't exist
        if not os.path.exists(example_file):
            prs = Presentation()
            # Slide 1: Title Slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            title.text = "Hello World from PPTX"
            subtitle.text = "python-pptx test presentation!"
            # Slide 2: Bullet Slide
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)
            title_shape = slide.shapes.title
            title_shape.text = "A Bullet Slide"
            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame
            tf.text = "Find Orders"
            p = tf.add_paragraph()
            p.text = "Find Customers"
            p.level = 1
            # Add a note to slide 2
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = "This is a note on slide 2."
            prs.core_properties.author = "PPTX Parser Test"
            prs.core_properties.title = "Sample PPTX"
            prs.save(example_file)
            # print("Created dummy example.pptx for testing.")
            logger.info("Created dummy example.pptx for testing.")
        
        # print(f"\n--- Parsed Data for {example_file} ---")
        logger.info(f"--- Parsed Data for {example_file} ---")
        parsed_data = parse_pptx(example_file)

        if parsed_data.get('error'):
            logger.error(f"Error parsing PPTX for test: {parsed_data['error']}")
        else:
            for key, value in parsed_data.items():
                if isinstance(value, list) and key in ["text_from_slides", "speaker_notes"]:
                    # print(f"\n{key.replace('_', ' ').title()}:")
                    logger.debug(f"{key.replace('_', ' ').title()}:")
                    for item in value[:2]: # print first 2 for brevity
                        # print(item)
                        logger.debug(str(item))
                    if len(value) > 2:
                        logger.debug("...")
                elif isinstance(value, dict) and key == "metadata":
                    # print(f"\n{key.title()}:")
                    logger.debug(f"{key.title()}:")
                    for meta_key, meta_val in value.items():
                        # print(f"  {meta_key.title()}: {meta_val}")
                        logger.debug(f"  {meta_key.title()}: {meta_val}")
                elif key != "error":
                    # print(f"\n{key.title()}: {value}")
                    logger.debug(f"{key.title()}: {value}")
            # print("-------------------------------------")
            logger.info("-------------------------------------")

    except ImportError:
        # print("python-pptx library is not installed. Skipping dummy PPTX creation and test.")
        logger.warning("python-pptx library is not installed. Skipping dummy PPTX creation and test.")
    except Exception as e:
        # print(f"An error occurred during the __main__ test block: {e}")
        logger.error(f"An error occurred during the __main__ test block: {e}", exc_info=True) 